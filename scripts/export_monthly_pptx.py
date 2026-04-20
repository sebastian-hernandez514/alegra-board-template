#!/usr/bin/env python3
"""
export_monthly_pptx.py — Genera 2 slides editables de Monthly/YTD Performance
desde data/metrics.yaml → output/monthly_performance.pptx

Usage:
    uv run --with python-pptx --with pyyaml python3 scripts/export_monthly_pptx.py
"""

from pathlib import Path
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
ROOT   = Path(__file__).parent.parent
OUTPUT = ROOT / "output" / "monthly_performance.pptx"

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x20, 0x35)
PURPLE = RGBColor(0x53, 0x4A, 0xB7)
GREEN  = RGBColor(0x1D, 0x9E, 0x75)
RED    = RGBColor(0xD8, 0x5A, 0x30)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x99, 0x99, 0x99)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG2    = RGBColor(0xF5, 0xF5, 0xF3)

# ── Helpers ───────────────────────────────────────────────────────────────────
def inches(v): return Inches(v)
def pt(v):     return Pt(v)

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, inches(x), inches(y), inches(w), inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=10, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True):
    txb = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Inter"
    return txb

def delta_color(positive, neutral=False):
    if neutral: return GRAY
    return GREEN if positive else RED

def add_kpi_block(slide, x, y, w,
                  label, value, value_size,
                  deltas,          # list of (label, value, positive)
                  accent=PURPLE):
    """Draws a KPI card: label · big value · delta pills"""
    # Label
    add_text(slide, label, x, y, w, 0.18, size=7, color=GRAY, bold=False)
    # Big value
    add_text(slide, value, x, y+0.17, w, 0.40, size=value_size, bold=True, color=DARK)
    # Deltas
    dx = x
    dy = y + 0.55
    pill_w = w / max(len(deltas), 1) - 0.04
    for (dlabel, dval, dpos) in deltas:
        bg = RGBColor(0xE8, 0xF5, 0xEE) if dpos else RGBColor(0xFC, 0xEE, 0xE9)
        tc = GREEN if dpos else RED
        if dval in ("N/A", "—", None, ""):
            bg = BG2; tc = GRAY
        r = add_rect(slide, dx, dy, pill_w, 0.18, fill=bg)
        add_text(slide, f"{dlabel} {dval}", dx+0.04, dy+0.02, pill_w-0.05, 0.15,
                 size=6.5, color=tc, bold=True)
        dx += pill_w + 0.05

def add_section_label(slide, text, x, y, w):
    add_rect(slide, x, y, w, 0.01, fill=PURPLE)
    add_text(slide, text.upper(), x, y+0.03, w, 0.16,
             size=6.5, bold=True, color=PURPLE)

# ── Load data ─────────────────────────────────────────────────────────────────
def load():
    with open(ROOT / "data" / "metrics.yaml") as f:
        return yaml.safe_load(f)

# ── Slide 1 — Monthly Performance ─────────────────────────────────────────────
def build_monthly(slide, d):
    W, H = 10.0, 5.63

    # Background
    add_rect(slide, 0, 0, W, H, fill=WHITE)

    # Left accent bar
    add_rect(slide, 0, 0, 0.08, H, fill=NAVY)

    # Header
    add_rect(slide, 0.08, 0, W-0.08, 0.65, fill=NAVY)
    add_text(slide, "Monthly Performance", 0.22, 0.05, 5, 0.28,
             size=16, bold=True, color=WHITE)
    add_text(slide, "February 2026", 0.22, 0.33, 5, 0.22,
             size=9, color=RGBColor(0xAA, 0xB4, 0xC8))
    add_text(slide, "Board Review · Alegra", W-2.5, 0.05, 2.3, 0.55,
             size=8, color=RGBColor(0x60, 0x70, 0x90), align=PP_ALIGN.RIGHT)

    # ── Section: Crecimiento ──────────────────────────────────────────────────
    add_section_label(slide, "Crecimiento", 0.22, 0.78, 9.6)

    kpi_y = 1.02
    col_w = 2.35

    # ARR
    add_kpi_block(slide, 0.22, kpi_y, col_w,
        "ARR", d["arr_total"], 20,
        [("MoM", d["arr_mom"], d["arr_mom_positive"]),
         ("YoY", d["arr_yoy"], d["arr_yoy_positive"]),
         ("vs Budget", d["arr_vs_budget"], d["arr_vs_budget_positive"])])

    # New MRR
    add_kpi_block(slide, 0.22 + col_w + 0.15, kpi_y, col_w,
        "New MRR", d["new_mrr"], 20,
        [("MoM", d["new_mrr_mom"], d.get("new_mrr_mom_positive", True)),
         ("YoY", d["new_mrr_yoy"], d.get("new_mrr_yoy_positive", True)),
         ("vs Budget", d["new_mrr_vs_budget"], d.get("new_mrr_vs_budget_positive", True))])

    # New Logos
    add_kpi_block(slide, 0.22 + (col_w + 0.15)*2, kpi_y, col_w,
        "New Logos", d["new_logos"], 20,
        [("MoM", d["new_logos_mom"], d.get("new_logos_mom_positive", True)),
         ("YoY", d["new_logos_yoy"], d.get("new_logos_yoy_positive", True)),
         ("vs Budget", d["new_logos_vs_budget"], d.get("new_logos_vs_budget_positive", False))])

    # Logo Churn
    churn_val = f"{d.get('logo_churn_global','?')}%"
    add_kpi_block(slide, 0.22 + (col_w + 0.15)*3, kpi_y, col_w,
        "Logo Churn (global)", churn_val, 20,
        [("Core", f"{d.get('logo_churn_core','?')}%", False),
         ("vs Budget", d.get("logo_churn_vs_budget_pp","—"), False)])

    # ── Section: Desempeño financiero ─────────────────────────────────────────
    add_section_label(slide, "Desempeño financiero", 0.22, 2.18, 9.6)

    kpi_y2 = 2.42
    col_w2 = 3.05

    # Net Revenue
    add_kpi_block(slide, 0.22, kpi_y2, col_w2,
        "Net Revenue", d["net_revenue"], 20,
        [("MoM", d["net_revenue_mom"], d["net_revenue_mom_positive"]),
         ("YoY", d["net_revenue_yoy"], d["net_revenue_yoy_positive"]),
         ("vs Budget", d["net_revenue_vs_budget"], d["net_revenue_vs_budget_positive"])])

    # Gross Margin
    add_kpi_block(slide, 0.22 + col_w2 + 0.15, kpi_y2, col_w2,
        "Gross Margin", d["gross_margin"], 20,
        [("MoM", d["gross_margin_mom"], d.get("gross_margin_mom_positive", True)),
         ("YoY", d["gross_margin_yoy"], d["gross_margin_yoy_positive"]),
         ("vs Budget", d["gross_margin_vs_budget"], d["gross_margin_vs_budget_positive"])])

    # EBITDA Margin
    add_kpi_block(slide, 0.22 + (col_w2 + 0.15)*2, kpi_y2, col_w2,
        "EBITDA Margin", d["ebitda_margin"], 20,
        [("MoM", d["ebitda_margin_mom"], d["ebitda_margin_mom_positive"]),
         ("YoY", d["ebitda_margin_yoy"], d["ebitda_margin_yoy_positive"]),
         ("vs Budget", d["ebitda_margin_vs_budget"], d["ebitda_margin_vs_budget_positive"])])

    # ── Section: Riesgo y sostenibilidad ──────────────────────────────────────
    add_section_label(slide, "Riesgo y sostenibilidad — Payback", 0.22, 3.58, 9.6)

    kpi_y3 = 3.82
    pb_col = 2.2

    add_kpi_block(slide, 0.22, kpi_y3, pb_col,
        "Payback Core (actual)", f"{d['payback_core']} mo", 18,
        [], accent=PURPLE)

    add_kpi_block(slide, 0.22 + pb_col + 0.15, kpi_y3, pb_col,
        "Payback Lite (actual)", f"{d['payback_lite']} mo", 18,
        [], accent=GREEN)

    add_kpi_block(slide, 0.22 + (pb_col + 0.15)*2, kpi_y3, pb_col,
        "Payback histórico (prom. 2025)", f"{d['payback_hist']} mo", 18,
        [], accent=GRAY)

    # Footer
    add_rect(slide, 0.08, H-0.22, W-0.08, 0.22, fill=BG2)
    add_text(slide, f"Fuente: Redshift · metrics.yaml · generado automáticamente",
             0.22, H-0.20, 6, 0.18, size=6, color=GRAY)
    add_text(slide, "Alegra Board Feb-26", W-2.2, H-0.20, 2.0, 0.18,
             size=6, color=GRAY, align=PP_ALIGN.RIGHT)

# ── Slide 2 — YTD Performance ─────────────────────────────────────────────────
def build_ytd(slide, d):
    W, H = 10.0, 5.63

    add_rect(slide, 0, 0, W, H, fill=WHITE)
    add_rect(slide, 0, 0, 0.08, H, fill=NAVY)
    add_rect(slide, 0.08, 0, W-0.08, 0.65, fill=NAVY)
    add_text(slide, "YTD Performance", 0.22, 0.05, 5, 0.28,
             size=16, bold=True, color=WHITE)
    add_text(slide, "Enero — Febrero 2026", 0.22, 0.33, 5, 0.22,
             size=9, color=RGBColor(0xAA, 0xB4, 0xC8))
    add_text(slide, "Board Review · Alegra", W-2.5, 0.05, 2.3, 0.55,
             size=8, color=RGBColor(0x60, 0x70, 0x90), align=PP_ALIGN.RIGHT)

    # ARR EoP
    add_section_label(slide, "ARR acumulado", 0.22, 0.78, 9.6)
    add_kpi_block(slide, 0.22, 1.02, 3.0,
        "ARR EoP (Feb-26)", d["arr_total"], 22,
        [("YoY", d["arr_yoy"], d["arr_yoy_positive"]),
         ("vs Budget", d["arr_vs_budget"], d["arr_vs_budget_positive"])])

    # New MRR YTD
    add_section_label(slide, "Crecimiento YTD", 0.22, 2.05, 9.6)

    col_w = 3.05
    add_kpi_block(slide, 0.22, 2.28, col_w,
        "New MRR YTD", d["new_mrr_ytd"], 20,
        [("YoY", d["new_mrr_ytd_yoy"], d["new_mrr_ytd_yoy_positive"])])

    add_kpi_block(slide, 0.22 + col_w + 0.15, 2.28, col_w,
        "New Logos YTD", d["new_logos_ytd"], 20,
        [("YoY", d["new_logos_ytd_yoy"], d["new_logos_ytd_yoy_positive"])])

    # Note
    add_rect(slide, 0.22, 3.55, 9.56, 0.50, fill=BG2)
    add_text(slide,
        "📌  YTD cubre Enero y Febrero 2026. Los vs Budget corresponden al período acumulado contra el plan anual prorrateado.",
        0.32, 3.60, 9.30, 0.40, size=8, color=GRAY)

    # Footer
    add_rect(slide, 0.08, H-0.22, W-0.08, 0.22, fill=BG2)
    add_text(slide, "Fuente: Redshift · metrics.yaml · generado automáticamente",
             0.22, H-0.20, 6, 0.18, size=6, color=GRAY)
    add_text(slide, "Alegra Board Feb-26", W-2.2, H-0.20, 2.0, 0.18,
             size=6, color=GRAY, align=PP_ALIGN.RIGHT)

# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    d = load()

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)

    blank = prs.slide_layouts[6]  # Blank layout

    s1 = prs.slides.add_slide(blank)
    build_monthly(s1, d)
    print("  ✅ Slide 1 — Monthly Performance")

    s2 = prs.slides.add_slide(blank)
    build_ytd(s2, d)
    print("  ✅ Slide 2 — YTD Performance")

    prs.save(OUTPUT)
    print(f"\n  💾 Guardado en: {OUTPUT}")

if __name__ == "__main__":
    main()
