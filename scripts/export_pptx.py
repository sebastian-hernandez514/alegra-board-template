#!/usr/bin/env python3
"""
export_pptx.py — board_standalone.html → boards/<period>/board.pptx

Usage:
    uv run --with playwright --with python-pptx python3 scripts/export_pptx.py
    uv run --with playwright --with python-pptx python3 scripts/export_pptx.py --input output/board_standalone.html
    uv run --with playwright --with python-pptx python3 scripts/export_pptx.py --out output/board.pptx

Requires Playwright browsers installed:
    uv run --with playwright python3 -m playwright install chromium
"""

import argparse
from pathlib import Path
from io import BytesIO

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Emu

ROOT       = Path(__file__).parent.parent
INPUT_FILE = ROOT / "output" / "board_standalone.html"
OUTPUT_FILE = ROOT / "output" / "board.pptx"

# Slide dimensions match the HTML CSS vars (960×540 px at 96dpi → EMU)
SLIDE_W_PX = 960
SLIDE_H_PX = 540
SCALE      = 2          # 2× for retina-quality screenshots
PX_TO_EMU  = 914400 / 96


def px_to_emu(px):
    return int(px * PX_TO_EMU)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", default=str(INPUT_FILE))
    parser.add_argument("--out",   default=str(OUTPUT_FILE))
    args = parser.parse_args()

    input_path  = Path(args.input).resolve()
    output_path = Path(args.out)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"📄 Abriendo {input_path.name}…")

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page    = browser.new_page(viewport={"width": SLIDE_W_PX * SCALE,
                                              "height": SLIDE_H_PX * SCALE})
        page.goto(f"file://{input_path}", wait_until="domcontentloaded")
        page.wait_for_timeout(3000)  # allow charts/fonts to render

        slides = page.query_selector_all(".slide, .fp-slide, .gtm-slide, .hc-slide")
        print(f"   {len(slides)} slides encontrados")

        # Build PPTX
        prs = Presentation()
        prs.slide_width  = px_to_emu(SLIDE_W_PX)
        prs.slide_height = px_to_emu(SLIDE_H_PX)
        blank_layout = prs.slide_layouts[6]  # completely blank

        for i, slide_el in enumerate(slides, 1):
            img_bytes = slide_el.screenshot()  # auto-clips to element bounds
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                BytesIO(img_bytes),
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"   [{i:>3}/{len(slides)}] slide capturado", end="\r")

        browser.close()

    prs.save(output_path)
    print(f"\n✅ PPTX guardado → {output_path}")
    print(f"   Ábrelo en Google Slides: File → Import slides → sube el .pptx")


if __name__ == "__main__":
    main()
